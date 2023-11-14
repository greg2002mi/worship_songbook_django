from django.shortcuts import render, redirect, get_object_or_404, HttpResponseRedirect
from django.http import HttpResponse, JsonResponse
from django.conf import settings
from django.core.paginator import Paginator
from django.core.paginator import (EmptyPage, PageNotAnInteger,
Paginator)
from django.views import generic
from django.urls import reverse, reverse_lazy
from .models import Post, Mlinks, Tag, Song, Lists, ListItem, Image, LANG, CHORDNOTE, Audio, Issues, BgImg
from django.contrib.auth.models import User, Group
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.decorators import login_required
from django.contrib.auth.forms import AuthenticationForm, PasswordChangeForm
from django.contrib import messages
from django.contrib.messages.views import SuccessMessageMixin
from django.utils.translation import gettext as _
from .forms import (AddSongForm, Transpose, AddTagForm, PostForm, EditPostForm, TagsForm, ContactUsForm, 
                    SongsForm, EmptyForm, EditSongForm, Assign2Event, NewUForm, UpdateUForm, UpdateProfileForm, 
                    UploadAvatarForm, AddSongTagForm, AddMediaForm, AddEventForm)
from .core import Chordpro_html
from .slides import CreatePptx
from django import forms
import os, uuid, json, re

from uuid import uuid4


from django.db.models import Q
from datetime import datetime
from django.utils import timezone
from django.core.mail import send_mail  # later to send mail - send_mail('Subject here', 'Here is the message.', 'from@example.com', ['to@example.com'], fail_silently=False)
from django.contrib.auth.views import PasswordResetConfirmView, PasswordResetView


## pptx creation
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from io import BytesIO
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
## pptx creation

PAGE_SIZE = getattr(settings, "PAGE_SIZE", 50)

# splits text into two columns 
def split_text(text, control, ori_key_int, transpose):
    lyrics = ""
    lines = text.split('\n')
    first_part = []
    second_part = []
    count = 0
    half_total = len(lines)/2
    html1 = ""
    html2 = ""
    
    
    for line in lines:
        if count < half_total:
            first_part.append(line)
        else:
            if line.strip() == '':
                second_part.append(line)
                second_part.extend(lines[count+1:])
                break
            else:
                first_part.append(line)
        count += 1
    
    first = '\n'.join(first_part) 
    second = '\n'.join(second_part)
    
    if control == 1:
        html1 = Chordpro_html(first, 1, ori_key_int, transpose)
        html2 = Chordpro_html(second, 1, ori_key_int, transpose)
    elif control ==2:
        html1 = Chordpro_html(first, False, 0, 0)
        html2 = Chordpro_html(second, False, 0, 0)
    
    lyrics = '<div class="row"><div class="col-auto">&nbsp;</div><div class="col-5">{}</div><div class="col-5"><br><br><br><br><br>{}</div></div>'.format(html1, html2)
    
    return lyrics

# class CustomPasswordResetConfirmView(PasswordResetConfirmView):
#     template_name = 'registration/change_password_form.html'
#     post_reset_login = True
#     success_url = reverse_lazy('index')

class ResetPasswordView(SuccessMessageMixin, PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject'
    success_message = "We've emailed you instructions for setting your password, " \
                      "if an account exists with the email you entered. You should receive them shortly." \
                      " If you don't receive an email, " \
                      "please make sure you've entered the address you registered with, and check your spam folder."
    success_url = reverse_lazy('users-home')

def login_page(request):
    if request.method == "POST":
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            username = form.cleaned_data.get('username')
            password = form.cleaned_data.get('password')
            user = authenticate(username=username, password=password)
            if user is not None:
                login(request, user)
                messages.info(request, _("You are now logged in as {}.").format(username))
                return redirect("index")
            else:
                messages.error(request, _("Invalid username or password."))
        else:
            messages.error(request, _("Invalid username or password."))
    form = AuthenticationForm()
    return render(request, "login.html", context={"form":form})

def logout_page(request):
    logout(request)
    messages.info(request, _("You have successfully logged out."))
    return redirect('index')

def register_page(request):
    if request.method == "POST":
        form = NewUForm(request.POST)
        if form.is_valid():
            user = form.save()
            login(request, user)
            messages.success(request, _("Registration successful."))
            return redirect("login")
        messages.error(request, _("Unsuccessful registration. Invalid information."))
    form = NewUForm()
    return render (request, "register.html", context={"form":form})
    
@login_required
def change_password(request):
    user = request.user
    if request.method == 'POST':
        form = PasswordChangeForm(user, request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, "Your password has been changed")
            return redirect('login')
        else:
            for error in list(form.errors.values()):
                messages.error(request, error)

    form = PasswordChangeForm(user)
    return render(request, 'password_change_form.html', {'form': form})   

@login_required
def profile(request):
    avatar_form = UploadAvatarForm()
    me = request.user
    if request.method == 'POST':
        
        user_form = UpdateUForm(request.POST, instance=me)
        profile_form = UpdateProfileForm(request.POST, request.FILES, instance=me.profile)

        if user_form.is_valid() and profile_form.is_valid():
            user_form.save()
            profile_form.save()
            messages.success(request, 'Your profile is updated successfully')
            return redirect('profile')
    else:
        user_form = UpdateUForm(instance=me)
        profile_form = UpdateProfileForm(instance=me.profile)
        songlist = me.song.all()
        postlist = me.post.all()
    context = {
        'title': _('Profile'),
        'user_form': user_form,
        'profile_form': profile_form,
        'avatar_form': avatar_form,
        'songlist': songlist,
        'postlist': postlist,
    }
    
    return render(request, 'user.html', context)

# function to upload avatar, and delete previous if exists.
@login_required
def avatar_upload(request, user_id):
    user = get_object_or_404(User, pk=user_id)
    if request.method == 'POST':
        my_file = request.FILES.get('avatar')
        if my_file:
            if user.profile.avatar:
                old_avatar_path = user.profile.avatar.path
                if os.path.isfile(old_avatar_path):
                    os.remove(old_avatar_path)
        newname = '{}{}'.format(uuid.uuid4().hex, os.path.splitext(my_file.name)[1])
        my_file.name = newname # change name of image
        profile = user.profile
        profile.avatar=my_file
        profile.save()
        return redirect('profile')
    messages.error(request, _("Error, failed to POST."))
    return redirect('profile')


class PostList(generic.ListView):
    queryset = Post.objects.filter(status=1).order_by('-created_on')
    template_name = 'index.html'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['form'] = PostForm()  # Replace 'MyForm' with the name of your form class
        return context

def make_post(request):
    if request.method == 'POST':
        form = PostForm(request.POST)
        if form.is_valid(): 
            post = form.save(commit=False)
            post.user = request.user
            post.save()
            post_id = post.id
            messages.success(request, _("Your post is drafted."))
            return redirect('confirm_post', post_id)
        else:
            messages.error(request, _("Not validated."))
            return redirect('index')
        
def confirm_post(request, post_id):
    post = get_object_or_404(Post, pk=post_id)
    if request.method == 'POST':
        form = EditPostForm(request.POST, instance=post)
        if form.is_valid(): 
            form.save()
            messages.success(request, _("Your post is on air."))
            return redirect('index')
    else:
        context = {
            'form': EditPostForm(instance=post),
            'post': post, 
            }
        return render(request, 'confirm_post.html', context) 

def edit_post(request, post_id):
    post = get_object_or_404(Post, pk=post_id)
    if request.method == 'POST':
        form = EditPostForm(request.POST, instance=post)
        if form.is_valid(): 
            form.save()
            messages.success(request, _("Your post is on air."))
            return redirect('index')
    else:
        context = {
            'form': EditPostForm(instance=post),
            'post': post, 
            }
        return render(request, 'edit_post.html', context)


def post_detail(request, post_id):
    post = get_object_or_404(Post, pk=post_id)
    context = {
        'title': _('Details of post:{}'.format(post.title)),
        'post': post,
        }
    return render(request, 'post_detail.html', context)


@login_required
def explore(request):
    current_datetime = timezone.now()
    songs = Song.objects.filter(status=1).order_by('-timestamp')
    new_songs = songs[:40]
    events = Lists.objects.filter(date_time__gt=current_datetime).order_by('date_time')
    o_events = Lists.objects.filter(date_time__lt=current_datetime).order_by('date_time')
    context = {
        'title': _('Explore'),
        'new_songs': new_songs,
        'events': events,
        'o_events': o_events,
        # 'next_url': next_url,
        # 'prev_url': prev_url,
    }
    return render(request, 'explore.html', context)

@login_required
def songbook(request):  
    search_songs = request.GET.get('search')
    if search_songs:
        songs = Song.objects.filter(Q(title__icontains=search_songs) & Q(lyrics__icontains=search_songs)).order_by('title')
        songs_all = Song.objects.filter(Q(title__icontains=search_songs) & Q(lyrics__icontains=search_songs)).order_by('title')
    else:    
        songs = Song.objects.filter(status=1).order_by('title')
        songs_all = Song.objects.order_by('title')
    form = AddSongForm()
    tags = Tag.objects.all()
    # page = request.GET.get('page', 1)
    # songs_list = Song.objects.order_by('title')
    paginator = Paginator(songs, PAGE_SIZE)
    page_number = request.GET.get("page")
    try:
        page_obj = paginator.get_page(page_number)
        
    except PageNotAnInteger:
        page_obj = paginator.get_page(1)
    except EmptyPage:
        page_obj = paginator.get_page(paginator.num_pages)
    # songs = paginator.get_page(page)
    # next_url = reverse('songbook') + '?page=' + str(songs.next_page_number()) if songs.has_next() else None
    # prev_url = reverse('songbook') + '?page=' + str(songs.previous_page_number()) if songs.has_previous() else None
    
    context = {
        'title': _('Songbook'),
        'songs': songs,
        'songs_all': songs_all,
        # 'next_url': next_url,
        # 'prev_url': prev_url,
        'keyset': CHORDNOTE,
        'lang': LANG,
        'form': form,
        'tags': tags,
        'page_obj': page_obj,
        'c': 2, # here c means return to 2 - return to songbook. 
        'search_word': search_songs,
        # change later
    }
    return render(request, 'songbook.html', context)

@login_required
def add_song(request):
    if request.method == 'POST':
        form = AddSongForm(request.POST)
        if form.is_valid():
            song = form.save(commit=False)
            song.publisher = request.user
            song.save()
            messages.success(request, _("Song successfully created as a draft."))
            return redirect('view_song', song_id=song.pk, key=song.key)  # Redirect to the songbook page after successful submission
    else:
        form = AddSongForm()
    
    return render(request, 'add_song.html', {'form': form})

@login_required
def publish(request, song_id):
    song = get_object_or_404(Song, pk=song_id)
    song.status = 1
    song.save()
    messages.success(request, _("Song has successfully been published."))
    return redirect('songbook')

def view_song(request, song_id, key):
    song = get_object_or_404(Song, pk=song_id)
    tags = Tag.objects.all()
    tag_states = {tag.id: tag in song.tags.all() for tag in Tag.objects.all()}
    images = song.image.all()
    media = song.media.all()
    tagged_list = []
    if song.tags.exists():
        tagged_list = [tagged.name for tagged in song.tags.all()]
    # ori_key = ''
    ori_key_int = song.key
    lyrics = song.lyrics
    # upload files
    media_root = settings.MEDIA_ROOT
    # directory = os.path.join(media_root, 'photos')
    files = os.listdir(media_root)
    audio = song.audio.all()
    form = Transpose()
    remove_transl_form = EmptyForm()
    untagall_form = EmptyForm()
    #to get a list of songs with condition, other language, and not already linked ones
    # other_songs = Song.objects.filter(Song.language!=song.language).all()
    # if song.translated is not None:
    #     linked = song.translated.all()
    #     for l in linked:
    #         other_songs.remove(l)
    # tr_list = [(s.id, s.title) for s in other_songs]
    # transl_form.tr_song_id.choices = tr_list
    # list of translated songs
    t1_songlist = song.translated.order_by('title').all()
    t2_songlist = song.translation.order_by('title').all()
    t_songlist = list(t1_songlist)
    for song_obj in t2_songlist:
        if song_obj not in t_songlist:
            t_songlist.append(song_obj)
    # delete_form = EmptyForm()
    
    tags_form = TagsForm(choices=[(tag.id, tag.name) for tag in Tag.objects.all()])
    # populate choices for tags_form from db
    # choices = [(t.id, t.name) for t in tags]
    # tags_form.name.choices = choices
    # print(choices)
    if key is None:
        transpose = ori_key_int
    else: 
        transpose = int(key)
    # to send original key as string
    key_dict = {'key': transpose}
    ori_key = CHORDNOTE[ori_key_int][1]
    if transpose != ori_key_int:
        ori_key = ori_key + _(" | transposed to ") + CHORDNOTE[transpose][1]
       
    html = Chordpro_html(song.lyrics, True, song.key, transpose)
    only_lyrics = Chordpro_html(lyrics, False, 0, 0)

    if song.lyrics is None:
        messages.success(request, _('Current song does not have lyrics'))
        return redirect('songbook')
    else:
        
        if request.method == 'POST':
            tags_form = TagsForm(request.POST)
            form = Transpose(request.POST)
            if form.is_valid():
                # Handle the submit of transpose chord action
                key = form.cleaned_data['key']
                return redirect('view_song', song_id=song.id, key=key)
            # if tags_form.is_valid():
            #     selected_tags = tags_form.cleaned_data['name']
            #     song.tags.clear()  # Clear existing tags
            #     for tag_id in selected_tags:
            #         tag = Tag.objects.get(id=tag_id)
            #         song.tags.add(tag)
            #     song.save()
            #     # Flash message and redirect
            #     return redirect('view_song', song_id=song.id, key=song.key)
        elif request.method == 'GET':
            
            initial_data = {'name': [tag.id for tag in song.tags.all()]}
            tags_form = TagsForm(initial=initial_data)    
            
        context = {
            'title': _('View Song'),
            'song': song,
            'keyset': CHORDNOTE,
            'lang': LANG,
            'tags': tags,
            'html': html, 
            'transl_form': SongsForm(current_song=song), 
            'remove_transl_form': remove_transl_form,  
            'tagged_list': tagged_list, 
            'tag_states': tag_states, 
            'tags_form': tags_form, 
            'only_lyrics': only_lyrics, 
            'form': Transpose(initial=key_dict), 
            'ori_key': ori_key, 
            't_songlist': t_songlist,  
            'untagall_form': untagall_form, 
            'media': media,
            'images': images,
            'audio': audio,
            'files': files,
            
        }
        return render(request, 'view_song.html', context)

def download_presentation(request, song_id, bg):
    song = get_object_or_404(Song, pk=song_id)
    lyrics = song.lyrics
    # make bg model that stores various backgrounds for presentation.
    bgimg = get_object_or_404(BgImg, pk=song_id)
    # Create presentation
    presentation = CreatePptx(lyrics, bgimg.image.url)
    
    # Save the presentation to a BytesIO buffer
    from io import BytesIO
    buffer = BytesIO()
    presentation.save(buffer)
    
    # Create the HttpResponse object with the presentation content
    response = HttpResponse(buffer.getvalue(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    response["Content-Disposition"] = "attachment; filename={}.pptx".format(song.title)

    return response

@login_required
def delete_song(request, song_id):
    song = get_object_or_404(Song, pk=song_id)
    song.delete()
    # delete song from database
    messages.success(request, _('Song has been permanently deleted.'))
    return redirect('songbook')

@login_required
def edit_song(request, song_id):
    song = get_object_or_404(Song, pk=song_id)
    if request.method == 'POST':
        # to edit add instance of that song, otherwise a new entry will be saved instead
        form = EditSongForm(request.POST, instance=song)
        if form.is_valid():
            form.save()
            messages.success(request, _('Song has been updated successfully.'))
            return redirect('view_song', song_id=song_id, key=song.key)
    else:
        context = {
            'form': EditSongForm(instance=song),
            'song_id': song_id, 
            'song': song,
            }
        return render(request, 'edit_song.html', context)

@login_required
def upload_audio(request, song_id):
    song = get_object_or_404(Song, pk=song_id)
    if request.method == 'POST':
        my_file = request.FILES.get('file')
        media = song.audio.all()
        num = len(media) + 1
        newname = "{}.{}_{}.mp3".format(song.id, song.title, str(num))
        original_filename = my_file.name
        my_file.name = newname # change name of image
        audio = Audio.objects.create(audio_file=my_file, title=original_filename, song=song, filename=newname)
        audio.save()
        song.audio.add(audio)
        return redirect('manage_media', song_id=song.id)
    return '', 204

@login_required
def upload_video(request, song_id):
    song = get_object_or_404(Song, pk=song_id)
    if request.method == 'POST':
        form = AddMediaForm(request.POST)
        if form.is_valid():
            video = form.save(commit=False)
            video.save()
            song.media.add(video)
            messages.success(request, _('Video link has been successfully added.'))
            return redirect('manage_media', song_id=song.id)

@login_required
def file_upload(request, song_id, mtype):
    song = get_object_or_404(Song, pk=song_id)
    if request.method == 'POST':
        my_file = request.FILES.get('file')
        media = song.image.all()
        num = len(media) + 1
        newname = "{}.image_{}.jpg".format(song.id, str(num))
        original_filename = my_file.name
        my_file.name = newname # change name of image
        Image.objects.create(image=my_file, title=original_filename, song=song, filename=newname)
        context = {
            'song_id': song_id, 
            'song': song,
            'key': song.key,
            }
        return HttpResponse('view_song', context)
    return JsonResponse({'post':'false'})

def wrapper(instance, filename):
    ext = filename.split('.')[-1]
    # get filename
    if instance.pk:
        filename = '{}.{}'.format(instance.pk, ext) # do instance.username 
                                                    # if you want to save as username
    else:
        # set filename as random string
        filename = '{}.{}'.format(uuid4().hex, ext)
    # return the whole path to the file
    return os.path.join('images/', filename)

# def delete_file(request, song_id, mlink_id):
#     pass
@login_required
def manage_media(request, song_id):
    song = get_object_or_404(Song, pk=song_id)
    if song.image:
        images = song.image.all()
    else:
        images = []
    if song.audio:
        music = song.audio.all()
    else:
        music = []
    if song.media:
        media = song.media.all()
    else:
        media = []
    form = AddMediaForm()
    context = {
        'song': song, 
        'form': form,
        'images': images,
        'key': song.key,
        'music': music,
        'media': media,
        }
    return render(request, 'manage_media.html', context)

@login_required
def delete_image(request, song_id, image_id):
    song = get_object_or_404(Song, pk=song_id)
    image = get_object_or_404(Image, pk=image_id)
    if request.method == 'POST':
        # song.image.remove(image)
        image.delete()
        messages.success(request, _('Image has been successfully deleted.'))
        return redirect('manage_media', song_id=song.id)

@login_required
def delete_audio(request, song_id, audio_id):
    song = get_object_or_404(Song, pk=song_id)
    audio = get_object_or_404(Audio, pk=audio_id)
    if request.method == 'POST':
        # song.audio.remove(audio)
        audio.delete()
        messages.success(request, _('Audio file has been successfully deleted.'))
        return redirect('manage_media', song_id=song.id)

@login_required    
def delete_video(request, song_id, mlink_id):
    song = get_object_or_404(Song, pk=song_id)
    video = get_object_or_404(Mlinks, pk=mlink_id)
    if request.method == 'POST':
        song.media.remove(video)
        video.delete()
        messages.success(request, _('Audio file has been successfully deleted.'))
        return redirect('manage_media', song_id=song.id)

# def upload_i(request, filename):
#     pass

@login_required
def add_transl(request, song_id):
    current_song = get_object_or_404(Song, pk=song_id)
    if request.method == 'POST':
        transl_form = SongsForm(request.POST, current_song=current_song)
        if transl_form.is_valid():
            oth_song_id = transl_form.cleaned_data['title']
            other_song = get_object_or_404(Song, pk=oth_song_id)
            if current_song.language == other_song.language:
                messages.error(request, _('Only same songs in different languages can be bind.'))
                return redirect('view_song', song_id=current_song.id, key=current_song.key)
            current_song.translated.add(other_song)
            messages.success(request, _('Translation has been successfully linked.'))
            return redirect('view_song', song_id=current_song.id, key=current_song.key)

@login_required
def remove_transl(request, selsong_id, cursong_id):
    current_song = get_object_or_404(Song, pk=cursong_id)
    selected_song = get_object_or_404(Song, pk=selsong_id)
    if request.method == 'POST':
        if selected_song in current_song.translated.all():
            current_song.translated.remove(selected_song)
        elif current_song in selected_song.translated.all():
            selected_song.translated.remove(current_song)
            messages.success(request, _('Translation has been successfully unlinked.'))
            return redirect('view_song', song_id=current_song.id, key=current_song.key)
        else:
            messages.error(request, _('Error occured. Please reload page and try again.'))
            return redirect('view_song', song_id=current_song.id, key=current_song.key)


def tag_list(request):
    tags = Tag.objects.all()
    form = AddTagForm()
    deletetag_form = EmptyForm()
    if request.method == 'POST':
        form = AddTagForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, _('Tag is saved successfully.'))
            return redirect('tag_list')
    context = {
        'title': _('Tag list'), 
        'tags': tags, 
        'form': form, 
        'deletetag_form': deletetag_form,
        }
    return render(request, 'tag_list.html', context)

def tag_songlist(request, tag_id):
    tag = get_object_or_404(Tag, pk=tag_id)
    songs = tag.song.order_by('title')
    paginator = Paginator(songs, PAGE_SIZE)
    page_number = request.GET.get("page")
    addform = AddSongTagForm(tag=tag)
    
    try:
        page_obj = paginator.get_page(page_number)
        
    except PageNotAnInteger:
        page_obj = paginator.get_page(1)
    except EmptyPage:
        page_obj = paginator.get_page(paginator.num_pages)
    # make searchable drop down select
    # if request.method == 'POST':
    #     addform = AddSongTagForm(request.POST)
    #     if addform.is_valid():
    #         song = addform.cleaned_data.get('id')
    #         print(song)
    context = {
        'title': _('Songbook'),
        'songs': songs,
        'keyset': CHORDNOTE,
        'lang': LANG,
        'addform': addform,
        'tag': tag,
        'page_obj': page_obj,
    }
    
        
    return render(request, 'tag_songlist.html', context)

@login_required    
def delete_tag(request, tag_id):
    tag = get_object_or_404(Tag, pk=tag_id)
    if request.method == 'POST':
        tag.delete()
        messages.success(request, _('Tag is successfully deleted.'))
        return redirect('tag_list')
    else:
        messages.error(request, _('Error, failed to delete tag.'))
        return redirect('tag_list')


@login_required    
def tagging(request, song_id):
    song = get_object_or_404(Song, pk=song_id)
    # tag_states = {tag.id: tag in song.tags.all() for tag in Tag.objects.all()}
    tags_form = TagsForm(request.POST)
    
    if request.method == "POST":
        print(tags_form)
        if tags_form.is_valid():
            selected_tags = tags_form.cleaned_data['name'] 
            tagged_list = song.tags.all()
            # if not tagged_list:
            #     for dt in tagged_list:
            #         song.tags.remove(dt)
            
            #first check if unchecked
            for tg in tagged_list:
                if tg not in selected_tags:
                    song.tags.remove(tg)
        
            # add new tags
            for t in selected_tags:
                if t not in tagged_list:
                    song.tags.add(t)
            return redirect('view_song', song_id=song.id, key=song.key) 
        else:
            messages.error(request, _('Error, failed to validate.'))
            return redirect('view_song', song_id=song.id, key=song.key) 
    else:
        messages.error(request, _('Error, failed to add tag to this song.'))
        return redirect('view_song', song_id=song.id, key=song.key)    

@login_required    
def untag(request, tag_id, song_id):
    song = get_object_or_404(Song, pk=song_id)
    tag = get_object_or_404(Tag, pk=tag_id)
    song.tags.remove(tag)   
    messages.success(request, _('{} tag is no longer related to this song.'.format(tag.name)))
    return redirect('tag_songlist', tag_id=tag_id)
                    
@login_required                    
def untagall(request, song_id):
    song = get_object_or_404(Song, pk=song_id)
    for t in song.tags.all():
        song.tags.remove(t)
    song.save()
    messages.success(request, _('Current song now has no tags.'))
    return redirect('view_song', song_id=song.id, key=song.key)


# cart

@login_required
def cart(request):
    current_user = request.user
    group = Group.objects.get(name="Minister")
    users = group.user_set.all()
    songlist = current_user.cart.all().order_by('listorder')
    context = {
        'title': _('Add event'),
        'songlist': songlist,
        'keyset': CHORDNOTE,
        'users': users,
        }
    return render(request, 'cart.html', context)

@login_required
def delete_item(request, item):
    current_user = request.user 
    for i in current_user.cart.all():
        if i.id == item:
            current_user.cart.remove(i)
    return redirect('cart')        

@login_required
def empty_cart(request):
    current_user = request.user
    current_user.cart.clear()
    return redirect('songbook')

@login_required
def assign2event(request):
    addform = AddEventForm()
    assignform = Assign2Event()
    context = {
        'title': _('Assign to event'),
        'addform': addform,
        'assignform': assignform,
        }       
    return render(request, 'assign2event.html', context) 

@login_required
def add2event(request):
    current_user = request.user
    if request.method == "POST":
        addform = AddEventForm(request.POST)
        if addform.is_valid():
            event = addform.save(commit=False)
            event.user = current_user
            event.save()
            for item in current_user.cart.all():
                event.items.add(item)
            current_user.cart.clear()
            messages.success(request, _('Success! Event is set and list added.'))
            return redirect('calendar')
        else:
            messages.error(request, _('Error! Form failed to validate.'))
            return redirect('assign2event')

@login_required
def add2xevent(request):
    current_user = request.user
    if request.method == "POST":
        assignform = Assign2Event(request.POST)
        if assignform.is_valid():

            event_id = assignform.cleaned_data['title']
            event = get_object_or_404(Lists, pk=event_id)
            event_has = event.items.count()
            for item in current_user.cart.all():
                if event_has > 0:
                    initial = item.listorder
                    item.listorder = initial + event_has
                    item.save()
                event.items.add(item)
            current_user.cart.clear()
            messages.success(request, _('Success! Cart list added to event {}-{}.'.format(event.date_time, event.title)))
            return redirect('calendar')
        else:
            messages.error(request, _('Error! Form failed to validate.'))
            return redirect('assign2event')    

@login_required
def add_to_cart(request, songid, c, keyword):
    song = get_object_or_404(Song, pk=songid)
    current_user = request.user
    cart = current_user.cart.all().order_by('listorder')
    
    if cart:
        lastitem = cart.last()
        if lastitem:
            lastitem_num = lastitem.listorder
        else:
            lastitem_num = 0
    else:
        lastitem_num = 0
    cartitem = ListItem(title=song.title, desired_key=song.key, listorder=lastitem_num+1)
    cartitem.save()
    
    cartitem.song.add(song)
    current_user.cart.add(cartitem)
    
    # need to finish
    if c == 1:
        return redirect('songbook')
    if c == 2:
        return redirect('songbook') # search
    else:
        return redirect('songbook') # explore

@login_required
def unsign_from_cartitem(request, item_id, username, state):
    listitem = ListItem.objects.get(pk=item_id)
    user = get_object_or_404(User, username=username)
    if user:
        listitem.assigned.remove(user)
    if state == 1:
        lists = listitem.lists.first()
        return redirect('lists', list_id = lists.id)
    if state == 2:
        return redirect('cart')

@login_required
def cart_update_list_order(request):
    if request.method == 'POST':
        order = request.POST.getlist('order[]')
        state = request.POST.get('state')
        for index, item_id in enumerate(order, start=1):
            listitem = get_object_or_404(ListItem, pk=item_id)
            listitem.listorder = index
            listitem.save()
        # return JsonResponse({'message': 'List order updated successfully.'})
        # if state == 1:
        lists = listitem.lists.first()
        return redirect('lists', list_id = lists.id)
        # if state == 2:
        #    return redirect('cart')

@login_required
def cart_update_desired_key(request):
    if request.method == 'POST':
        item_id = request.POST.get('item_id')
        desired_key = int(request.POST.get('desired_key'))
        state = request.POST.get('state')
        list_item = get_object_or_404(ListItem, pk=item_id)
        list_item.desired_key = desired_key
        list_item.save()
        # return JsonResponse({'message': 'Desired key updated successfully.'})
        # if state == 1:
        lists = list_item.lists.first()
        return redirect('lists', list_id = lists.id)
        # if state == 2:
        #    return redirect('cart')

@login_required
def cart_update_notes(request):
    if request.method == 'POST':
        item_id = request.POST.get('item_id')
        notes = request.POST.get('notes')
        state = request.POST.get('state')
        list_item = get_object_or_404(ListItem, pk=item_id)
        list_item.notes = notes
        list_item.save()
        # return JsonResponse({'message': 'Notes updated successfully.'})
        # if state == 1:
        lists = list_item.lists.first()
        return redirect('lists', list_id = lists.id)
        # if state == 2:
        #     return redirect('cart')

@login_required
def cart_assign_user(request):
    if request.method == 'POST':
        item_id = request.POST.get('item_id')
        user_id = request.POST.get('user_id')
        state = request.POST.get('state')
        listitem = get_object_or_404(ListItem, pk=item_id)
        user = get_object_or_404(User, pk=user_id)
        if user:
            listitem.assigned.add(user)
            listitem.save()
        # if state == 1:
        lists = listitem.lists.first()
        return redirect('lists', list_id = lists.id)
        # if state == 2:
        #     return redirect('cart')





# calendar
def calendar(request):
    # form = EmptyForm() # to delete events
    events = Lists.objects.order_by("created").all()
    for e in events:
        date = timezone.now()
        if e.date_end and e.date_end < date:  # ISSUE: Nonetype and datetime are not comparable
            e.status = 2
            e.save(update_fields=['status'])
    active_events = Lists.objects.filter(~Q(status=2))
    # check if date is passed
        
    context = {
        'title': _('Songbook'),
        'active_events': active_events,
        'lang': LANG,
    }
    return render(request, 'calendar.html', context)

def events(request):    
    events = []
    lists = Lists.objects.all()
    for list_item in lists:
        events.append({
            'id': list_item.id,
            'title': list_item.title,
            'start': list_item.date_time.strftime("%Y/%m/%d, %H:%M:%S"),
            'end': list_item.date_end.strftime("%Y/%m/%d, %H:%M:%S"),
            'url': '/lists/{}'.format(list_item.id),  # Link to the detail page
            'editable': True  # Set to False if you want to disable event editing
        })
        # event = {
        #     'id': list_item.id,
        #     'title': list_item.title,
        #     'start': list_item.date_time.isoformat() if list_item.date_time else None,  #datetime.datetime.strptime(str(i.start_date.date()), "%Y-%m-%d").strftime("%Y-%m-%d")
        #     'end': list_item.date_end.isoformat() if list_item.date_end else None,
        #     'url': '/lists/{}'.format(list_item.id),  # Link to the detail page
        #     'editable': True  # Set to False if you want to disable event editing
        # }
        # events.append(event)
    # return jsonify(events)
    return JsonResponse(events, safe=False)

@login_required
def add_event(request):
    form = AddEventForm()
    if request.method == "POST":
        form = AddEventForm(request.POST)
        if form.is_valid():
            event = form.save(commit=False)
            event.user = request.user
            event.save()
            messages.success(request, _('Success! Event is set.'))
            return redirect('calendar')
        else:
            messages.error(request, _('Error! Form failed to validate.'))
            return redirect('add_event')
    context = {
        'title': _('Add event'),
        'form': form,
        }
    return render(request, 'add_event.html', context)

@login_required
def edit_event(request, listid):
    event = get_object_or_404(Lists, pk=listid)
    if request.method == 'POST':
        # to edit add instance of that song, otherwise a new entry will be saved instead
        form = AddEventForm(request.POST, instance=event)
        if form.is_valid():
            form.save()
            messages.success(request, _('Event has been updated successfully.'))
            return redirect('lists', list_id=listid)
    else:
        context = {
            'form': AddEventForm(instance=event),
            'event': event
            }
        return render(request, 'edit_event.html', context)

@login_required
def delete_event(request, eventid, jump):
    event = get_object_or_404(Lists, pk=eventid)
    if request.method == 'POST':
        event.delete()
        messages.success(request, _('Event is removed successfully.'))
        # jump to redirect based on where it has been deleted
        if jump == 1:
            return redirect('calendar')

@login_required
def unsign_from_listitem(request, item_id, username, state):
    user = get_object_or_404(User, username=username)
    item = get_object_or_404(ListItem, pk=item_id)
    item.assigned.remove(user)
    listid = item.lists.first()
    return redirect('lists', list_id=listid.id)
    
@login_required
def list_delete_item(request, item, listid):
    event = get_object_or_404(Lists, pk=listid)
    for i in event.items.all():
        if i.id == item:
            event.items.remove(i)
    return redirect('lists', list_id=listid)


def jall_events(request):                                                                                                 
    all_events = Lists.objects.all()                                                                                    
    out = []                                                                                                             
    for event in all_events:                                                                                             
        out.append({                                                                                                     
            'id': event.id,
            'title': event.title,
            'start': event.date_time.strftime("%Y/%m/%d, %H:%M:%S"),
            'end': event.date_end.strftime("%Y/%m/%d, %H:%M:%S"),
            'url': '/lists/{}'.format(event.id),  # Link to the detail page
            'editable': True                                                            
        })                                                                                                               
                                                                                                                     
    return JsonResponse(out, safe=False)  

@login_required
def jadd_event(request):
    start = request.GET.get("start", None)
    end = request.GET.get("end", None)
    title = request.GET.get("title", None)
    event = Lists(title=str(title), date_time=start, date_end=end)
    event.save()
    data = {}
    return JsonResponse(data)

@login_required
def jupdate(request):
    start = request.GET.get("start", None)
    end = request.GET.get("end", None)
    title = request.GET.get("title", None)
    id = request.GET.get("id", None)
    event = Lists.objects.get(pk=id)
    event.date_time = start
    event.date_end = end
    event.title = title
    event.save()
    data = {}
    return JsonResponse(data)

@login_required
def jremove(request):
    id = request.GET.get("id", None)
    event = Lists.objects.get(pk=id)
    event.delete()
    data = {}
    return JsonResponse(data)

def lists(request, list_id):
    event = get_object_or_404(Lists, pk=list_id)
    de_form = EmptyForm()
    unroleform = EmptyForm()
    deleteform = EmptyForm()
    songlist = event.items.order_by('listorder').all()
    group = Group.objects.get(name="Minister")
    users = group.user_set.all()
    context = {
        'title': _('Event Details'),
        'event': event,
        'songlist': songlist,
        'keyset': CHORDNOTE,
        'users': users,
        'de_form': de_form,
        'unroleform': unroleform,
        'deleteform': deleteform,
        }
    return render(request, 'event_detail.html', context)

### presentation creation

# def createpresentation(request, list_id):
#     event = get_object_or_404(Lists, pk=list_id)

def create_list_items(song_id):
    song = get_object_or_404(Song, pk=song_id)
    lyrics = song.lyrics
    lyrics = lyrics.replace('\n', '\x0A')
    # Remove {Intro} to {Chorus} or {Verse} and clean up []
    cleaned_lyrics = re.sub(r'\{Intro\}.*?\{(?:Chorus|Verse \d+)\}', '', lyrics, flags=re.DOTALL)
    cleaned_lyrics = re.sub(r'\{Instrumental\}.*?\{(?:Chorus|Verse \d+|Pre-Chorus|Ending|Bridge)\}', '', cleaned_lyrics, flags=re.DOTALL)
    cleaned_lyrics = re.sub(r'\[.*?\]', '', cleaned_lyrics)
    # Remove all ... .. ..  .  .  extra dots and spaces
    cleaned_lyrics = re.sub(r'\.(?![a-zA-Z\n])\s*', '', cleaned_lyrics)

    # Split by section
    sections = re.split(r'\{(?:Chorus|Verse \d+|Pre-Chorus|Ending|Bridge)\}', cleaned_lyrics)
    # sections = [section.strip() for section in sections if section.strip()]

    # Create ListItem instances
    li = []
    for section in sections:
        lines = section.split('\n')
        li.append([line.strip() for line in lines if line.strip()])

    return li



def create_presentation(request, list_id, background_color=None):
    try:
        event = Lists.objects.get(pk=list_id)
        list_items = event.items.order_by("listorder").all()

        # Create a presentation
        prs = Presentation()
        
        slide_width = Inches(16)  # Adjust the width as needed (or use Cm for centimeters)
        slide_height = Inches(9)  # Adjust the height as needed (or use Cm for centimeters)
        prs.slide_width = slide_width
        prs.slide_height = slide_height
        
        for item in list_items:
            
            song = item.song.first()
            lyrlist = []
            lyrlist = create_list_items(song.id)
            
            if not len(lyrlist):
                slide = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[6])
                if background_color:
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = int(background_color, 16)
                else:
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
            else:
                for sl in lyrlist:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use a blank slide layout
        
                    # Set the background color or image for the slide
                    if background_color:
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = int(background_color, 16)
                    else:
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
        
                    # Add title and content to the slide
                    # shapes = slide.shapes
                    # title_shape = shapes.title
                    # title_shape.text = "Song " + str(item.listorder)
            # Add lyrics or other content to the slide
            # lyrics = " ".join([song.lyrics for song in item.song.all()])
                    # first textbox with a label of a song        
                    left = Cm(1)
                    top = Cm(1)
                    width = Inches(15)
                    height = Inches(1)
                    tx1Box = slide.shapes.add_textbox(left, top, width, height)
                    tf1 = tx1Box.text_frame
                    tf1.vertical_anchor = MSO_ANCHOR.TOP
                    tf1.word_wrap = False
                    tf1.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    p1 = tf1.add_paragraph()
                    p1.text = "Song " + str(item.listorder) + ": " + item.title
                    p1.font.bold = True
                    p1.font.size = Pt(12)
                    p1.font.color.rgb = RGBColor(105, 105, 105)  # White text
                    
                    # main textbox with lyrics
                    left = Cm(1)
                    top = Cm(3)
                    width = Inches(15)
                    height = Inches(7)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    tf.word_wrap = True
                    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    for sll in sl:
                        p = tf.add_paragraph()
                        p.text = sll
                        p.font.bold = True
                        p.font.size = Pt(36)
                        p.alignment = PP_ALIGN.CENTER
                        p.font.color.rgb = RGBColor(255, 255, 255)  # White text

        # Create a byte stream to save the presentation
        presentation_stream = BytesIO()
        prs.save(presentation_stream)
        presentation_stream.seek(0)

        # Create an HTTP response with the presentation
        response = HttpResponse(content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        response["Content-Disposition"] = f'attachment; filename="{event.title}.pptx"'
        response.write(presentation_stream.read())

        return response

    except Lists.DoesNotExist:
        # Handle the case where the list doesn't exist
        messages.error(request, _('Error, list is empty.'))
        return redirect('list', list_id=list_id)    
    

### testing presentation creation

# on stage mode. in general three modes. lyrics with chords, lyrics, or images
def onstage(request, eventid, viewtype):
    event = get_object_or_404(Lists, pk=eventid)

    unsorted = sorted(event.items.all(), key=lambda x: x.listorder)
    songlist = [item for item in unsorted]
    transpose = 0
    stage_mode = {}
    chords = []
    lyrics = []
    # images = []
    index_list = []
    image_state = [] # list of True false states if song has image
    music_sh = []
    if not unsorted:
        messages.error(request, _('Error, there are no items in this event.'))
        return redirect('lists', list_id=eventid)
    if not viewtype:
        messages.error(request, _('Error, view condition is not defined.'))
        return redirect('lists', list_id=eventid)
    # create two lists of lyrics with chords and without
    for i in unsorted:
        transpose = i.desired_key
        for s in i.song.all():
            ori_key_int = s.key
            # if song does not have chords lyrics without chords must be added.
            if s.key:
                split1 = split_text(s.lyrics, 1, ori_key_int, transpose) 
            else:
                split1 = split_text(s.lyrics, 2, 0, 0)    
            split2 = split_text(s.lyrics, 2, 0, 0)
        
            lyrics.append(split2)
            chords.append(split1)
            # check if song has images and append True state into list. otherwise False
            cntrl = s.image.count()
            
            if cntrl < 1:
                image_state.append(False)
            else:
                image_state.append(True)
            
    
    # create index_list and list of links of images
               
    # make a list of murls of pictures in a list
    for index, i in enumerate(unsorted):
        index_list.append(index + 1)
        for ss in i.song.all():
            # create list of lists of images
            if image_state[index]:
                music_patch = []
                for i in ss.image.all():
                    music_patch.append(i.image.url)    
                music_sh.append(music_patch)
            # if no image available append chords with lyrics
            elif not image_state[index]:
                ori_key_int = ss.key
                transpose = i.desired_key
                split = split_text(ss.lyrics, 1, ori_key_int, transpose)
                music_sh.append(split) 
    
    stage_mode = {
    "id": index_list,
    "state": image_state,
    "chords": chords,
    "lyrics": lyrics,
    "images": music_sh,
}
    
    stage_one = zip(stage_mode["id"], stage_mode["chords"], stage_mode["lyrics"])
    stage_two = zip(stage_mode["id"], stage_mode["state"], stage_mode["chords"], stage_mode["images"])
    context = {
        'title': _('On Stage'),
        'songlist': songlist, 
        'event': event,
        'lyrics': lyrics,
        'viewtype': viewtype, 
        'stage_mode': stage_mode,
        'stage_one': stage_one,
        'stage_two': stage_two,
        }
    return render(request, 'onstage.html', context)

# feedback 
def contact_us(request):
    issues = Issues.objects.order_by('timestamp', 'status').all()
    form = ContactUsForm()
    if request.method == "POST":
        form = ContactUsForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, _('Thank you for your feedback. If you have given us your contact information, we will get back to you as soon as possible'))
            return redirect('index')
        else:
            messages.error(request, _('Error! Form failed to validate.'))
            return redirect(add_event)
    context = {
        'title': _('Contact us'),
        'form': form,
        'issues': issues,
        }
    return render(request, 'contact_us.html', context)