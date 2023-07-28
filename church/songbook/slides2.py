import re
from pptx import Presentation
from pptx.util import Inches

def create_song_presentation(lyrics, output_filename, bg_path):
    lines = lyrics.split('\n')
    text = ""
    slides = []
    current_slide = None
    
    
    for line in lines:
        if line.startswith("{Intro}"):
            continue

        if line.startswith("{"):
            if current_slide:
                slides.append(current_slide)
            current_slide = Presentation().slides.add_slide(presentation.slide_layouts[6])
            textbox = current_slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
            frame = textbox.text_frame
            frame.word_wrap = True
            p = frame.add_paragraph()
            text = ""
            continue

        cleaned_line = re.sub(r'\[[^\]]+\]', '', line)
        if cleaned_line.strip():
            text += cleaned_line + "\n"

    if current_slide:
        slides.append(current_slide)

    for i, slide in enumerate(slides):
        textbox = slide.shapes[0].text_frame
        p = textbox.paragraphs[0]
        p.text = text.strip()

    if bg_path:
        for slide in slides:
            slide.background['background'].fill.solid()
            slide.background['background'].fill.fore_color.rgb = RGBColor(0, 0, 0)
            slide.background['background'].fill.fore_color.brightness = 0.5
            slide.background['background'].fill.transparency = 0.5

    presentation.save(output_filename)