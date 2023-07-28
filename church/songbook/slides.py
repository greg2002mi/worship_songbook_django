from pptx import Presentation
from pptx.util import Cm, Inches
import re

def clean_text(text):
    # Remove chords in square brackets from the text
    return re.sub(r'\[[^\]]+\]', '', text)

def CreatePptx(lyrics, bg_path):
    # Split the lyrics into lines
    lines = lyrics.split('\n')

    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Set custom slide width and height (inches or centimeters)
    slide_width = Inches(16)  # Adjust the width as needed (or use Cm for centimeters)
    slide_height = Inches(9)  # Adjust the height as needed (or use Cm for centimeters)
    presentation.slide_width = slide_width
    presentation.slide_height = slide_height

    # Define the margins in centimeters
    left_margin = Cm(2)
    top_margin = Cm(2)

    # Set the background image for the slides
    presentation.slide_master.background.fill.solid()
    presentation.slide_master.background.fill.fore_color.rgb = (255, 255, 255)  # Set the background color (white in this case)
    background = presentation.slide_master.background
    background_picture = background.background_picture
    background_picture.left = 0
    background_picture.top = 0
    background_picture.width = slide_width
    background_picture.height = slide_height
    background_picture.crop_left = 0
    background_picture.crop_right = 0
    background_picture.crop_top = 0
    background_picture.crop_bottom = 0
    background_picture.nudge_left = 0
    background_picture.nudge_top = 0

    slide = None
    text_box = None
    text = ""

    # Iterate through each line
    for line in lines:
        line = line.strip()
        if line.startswith("{Intro}"):
            continue
        if line.startswith("{") and line.endswith("}"):
            # If the line is a section header like {Intro}, {Chorus}, {Verse 1}, etc.
            if slide is not None and text:
                # Add the previous section's text to the slide
                text_frame = text_box.text_frame
                text_frame.text = text

                # Auto-size the text to fit the text box
                text_frame.auto_size = True

            # Create a new slide for the section
            slide_layout = presentation.slide_layouts[6]
            slide = presentation.slides.add_slide(slide_layout)

            # Set the slide background to the background image
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = (255, 255, 255)  # Set the background color (white in this case)
            slide.shapes.add_picture(bg_path, 0, 0, width=slide_width, height=slide_height)

            # Create a text box on the slide
            text_box = slide.shapes.add_textbox(left_margin, top_margin, slide_width - 2 * left_margin, slide_height - 2 * top_margin)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            # Reset the text variable for the new section
            text = ""
        elif slide is not None:
            # Append the line text to the section's text
            text += clean_text(line) + '\n'

    # Add the text of the last section to the last slide
    if slide is not None and text:
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.auto_size = True

    # Save the presentation to the specified output filename
    # presentation.save(output_filename)
    return presentation